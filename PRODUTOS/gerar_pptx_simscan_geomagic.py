#!/usr/bin/env python3
"""
SIMSCAN-E + Geomagic Design X — Pitch de Vendas
Enterfix Indústria, Comércio e Serviços Ltda.

Script gerador de apresentação PowerPoint — v1.0
Mantém slides 1-6 (empresa) e constrói novos slides de produto.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import shutil, os

# ─── Paths ────────────────────────────────────────────────────
BASE = r'c:\Users\paulo\OneDrive\Documentos\CODIGUINHO\PRODUTOS'
SRC  = os.path.join(BASE, r'devok_simse_images\Apresentação - Enterfix (2).pptx')
DST  = os.path.join(BASE, r'devok_simse_images\SIMSCAN-E + Geomagic - Enterfix.pptx')
IMG  = os.path.join(BASE, r'devok_simse_images')

# ─── Color Palette (same as original) ─────────────────────────
NAVY  = RGBColor(0x10, 0x2F, 0x54)   # primary dark navy
NAVY2 = RGBColor(0x14, 0x26, 0x3C)   # slightly lighter navy
BLUE  = RGBColor(0x0F, 0x94, 0xD5)   # accent blue
GREEN = RGBColor(0xA0, 0xBF, 0x30)   # enterfix green
DARK  = RGBColor(0x23, 0x29, 0x32)   # near-black body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EEE   = RGBColor(0xEE, 0xEE, 0xEE)   # light text on dark bg
MGRAY = RGBColor(0x40, 0x4B, 0x59)   # medium gray
LGRAY = RGBColor(0xF4, 0xF6, 0xF9)   # light background
DBOX  = RGBColor(0x08, 0x1E, 0x36)   # very dark box fill
GBOX  = RGBColor(0x1E, 0x3A, 0x1E)   # dark green box

# ─── Setup ────────────────────────────────────────────────────
print("=== SIMSCAN-E + Geomagic Design X — Gerador de PPTX ===")
print(f"Copiando base: {os.path.basename(SRC)}")
shutil.copy2(SRC, DST)

prs  = Presentation(DST)
W    = prs.slide_width    # 9144000 EMU = 20"
H    = prs.slide_height   # 5143500 EMU = 11.25"

print(f"Dimensões: {W/914400:.1f}\" x {H/914400:.1f}\"")
print(f"Slides no original: {len(list(prs.slides))}")

# ─── Utility Functions ─────────────────────────────────────────
def I(inches):
    return Inches(inches)

def add_rect(slide, left, top, width, height, fill_rgb=None):
    shp = slide.shapes.add_shape(1, I(left), I(top), I(width), I(height))
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp

def add_pic(slide, filename, left, top, width, height):
    """Add picture - JPG/PNG only (webp skipped with fallback)."""
    path = os.path.join(IMG, filename) if not os.path.isabs(filename) else filename
    if not os.path.exists(path):
        print(f"  ⚠ Imagem não encontrada: {filename}")
        return None
    if filename.lower().endswith('.webp'):
        print(f"  ℹ Convertendo webp para PNG: {filename}")
        try:
            from PIL import Image as PILImage
            im = PILImage.open(path)
            png_path = path.replace('.webp', '_conv.png')
            im.save(png_path, 'PNG')
            path = png_path
        except Exception as e:
            print(f"  ⚠ Não foi possível converter webp: {e}")
            return None
    try:
        return slide.shapes.add_picture(path, I(left), I(top), I(width), I(height))
    except Exception as e:
        print(f"  ⚠ Erro ao adicionar imagem {filename}: {e}")
        return None

def send_back(slide, *shapes):
    """Send shapes to the back of the z-order."""
    sp_tree = slide.shapes._spTree
    for shp in reversed(shapes):
        if shp is None:
            continue
        el = shp.element
        sp_tree.remove(el)
        sp_tree.insert(2, el)

def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False,
             wrap=True, name='Calibri'):
    tb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text          = text
        run.font.size     = Pt(size)
        run.font.bold     = bold
        run.font.italic   = italic
        run.font.color.rgb = color
        run.font.name     = name
    return tb

def add_lines(slide, left, top, width, height, items, name='Calibri'):
    """
    items: list of (text, size, bold, color, align [, space_before_pt])
    """
    tb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, size, bold, color, align = item[:5]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if len(item) > 5 and item[5]:
            from pptx.util import Pt as _Pt
            p.space_before = _Pt(item[5])
        if text:
            run = p.add_run()
            run.text           = text
            run.font.size      = Pt(size)
            run.font.bold      = bold
            run.font.color.rgb = color
            run.font.name      = name
    return tb

# ─── Delete product slides 7-21 (keep company slides 1-6) ─────
print("\nRemovendo slides de produto originais (7–21)...")
sldIdLst = prs.slides._sldIdLst
total    = len(list(prs.slides))

for i in range(total - 1, 5, -1):   # remove from back, keep indices 0-5
    el  = sldIdLst[i]
    rId = el.get(qn('r:id'))
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(el)

print(f"Slides restantes (empresa): {len(list(prs.slides))}")

blank = prs.slide_layouts[6]   # Blank layout

def NS():
    """Add new blank slide."""
    return prs.slides.add_slide(blank)

# ─────────────────────────────────────────────────────────────
# ════ SLIDES DE PRODUTO: SIMSCAN-E + GEOMAGIC DESIGN X ═════
# ─────────────────────────────────────────────────────────────

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 7 — SIMSCAN-E: Visão Geral                       ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 7: SIMSCAN-E — Visão Geral...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)
img = add_pic(sl, 'simse_feat2_bg.jpg', 9.2, 0, 10.8, 11.25)
dim = add_rect(sl, 9.2, 0, 10.8, 11.25, fill_rgb=NAVY2)

add_text(sl, 1, 0.9,  8, 0.7, 'SCANNER 3D INDUSTRIAL',
         size=16, bold=True, color=GREEN)
add_text(sl, 1, 1.7,  8, 2.0, 'SIMSCAN-E',
         size=80, bold=True, color=WHITE)
add_text(sl, 1, 3.8,  8, 0.7, 'Inteligente. Sem Fio. Tamanho de Bolso.',
         size=21, color=BLUE)
add_text(sl, 1, 4.65, 8, 2.2,
         'O SIMSCAN-E é um scanner 3D industrial portátil de nova geração com '
         'Wi-Fi 6 nativo, 81 linhas de laser azul e computação de borda embarcada. '
         'Precisão de até 0,020 mm e 6,3 milhões de medições por segundo em apenas 600 g.',
         size=15, color=EEE)

stats = [
    ('0,020\nmm',  'Precisão',   1.0),
    ('6,3M\npts/s','Medições/s', 3.3),
    ('600 g',      'Peso',       5.6),
    ('Wi-Fi 6',    'Sem fio',    7.9),
]
for val, lbl, x in stats:
    add_rect(sl, x, 7.2, 2.0, 3.0, fill_rgb=DBOX)
    add_text(sl, x+0.1, 7.4,  1.8, 1.8, val,
             size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, x+0.1, 9.3,  1.8, 0.6, lbl,
             size=11, color=MGRAY, align=PP_ALIGN.CENTER)

send_back(sl, dim, img, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 8 — Especificações de Performance                ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 8: Especificações de Performance...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=DARK)
img = add_pic(sl, 'simse_hl_laser81.jpg', 10.5, 1.5, 9.0, 9.3)
dim = add_rect(sl, 10.5, 0, 9.5, 11.25, fill_rgb=DBOX)

add_text(sl, 1.0, 0.7, 9, 0.7,  'ESPECIFICAÇÕES DE',
         size=18, bold=True, color=BLUE)
add_text(sl, 1.0, 1.4, 9, 1.5,  'PERFORMANCE',
         size=60, bold=True, color=WHITE)

specs = [
    ('81 Linhas de Laser Azul',
     'Triple Cross Technology — máxima cobertura e densidade por varredura', BLUE),
    ('Precisão de até 0,020 mm',
     'Acreditado ISO 17025 · VDI/VDE 2634 Parte 3 · JJF 1951', GREEN),
    ('6.300.000 Medições / segundo',
     '180 FPS — rastreamento fluido de superfícies complexas', BLUE),
    ('Wi-Fi 6 Nativo (802.11a/b/g/n/ac)',
     'Sem roteador externo — computação de borda integrada ao scanner', GREEN),
    ('600 g · 203 × 80 × 44 mm',
     'Uma mão · Baterias hot-swap · Operação –10 °C a 40 °C · Laser Classe II', BLUE),
]
y = 3.4
for title, desc, accent in specs:
    add_rect(sl, 1.0, y, 0.18, 0.5, fill_rgb=accent)
    add_text(sl, 1.4, y,       8.8, 0.55, title, size=16, bold=True, color=WHITE)
    add_text(sl, 1.4, y+0.55,  8.8, 0.5,  desc,  size=12, color=MGRAY)
    y += 1.42

send_back(sl, dim, img, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 9 — 3 Modos de Operação                          ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 9: 3 Modos de Operação...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=LGRAY)
add_text(sl, 0.5, 0.5, 19, 1.0, '3 MODOS DE OPERAÇÃO VERSÁTEIS',
         size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, 0.5, 1.55, 19, 0.6,
         'Selecione o modo ideal para cada aplicação, geometria e nível de detalhe',
         size=15, color=MGRAY, align=PP_ALIGN.CENTER)

modes = [
    ('simse_modo_ultrarapido.webp', BLUE,  'MODO ULTRARRÁPIDO',
     '63 linhas de laser — 6.300.000 pts/s',
     ['• 180 FPS — captura fluida e contínua',
      '• Área de varredura: 700 × 600 mm',
      '• Peças de médio e grande porte',
      '• Alta produtividade em shopfloor',
      '• Ideal para inspeção dimensional']),
    ('simse_modo_hyperfine.webp',  NAVY2, 'MODO HYPERFINE',
     '17 linhas paralelas de alta definição',
     ['• Resolução máxima: 0,020 mm',
      '• Geometrias complexas e orgânicas',
      '• Componentes aeroespaciais críticos',
      '• Bordas afiadas e chanfros finos',
      '• Superfícies espelhadas e polidas']),
    ('simse_modo_furos.webp',      GREEN, 'MODO FUROS PROFUNDOS',
     '1 linha extra dedicada a cavidades',
     ['• Furos, cavidades e profundidades',
      '• Câmera de curta distância',
      '• Ângulo acentuado de captura',
      '• Espaços internos inacessíveis',
      '• Medição de posição de furos']),
]

col_w = 5.8
for i, (img_name, color, title, subtitle, bullets) in enumerate(modes):
    x = 0.5 + i * (col_w + 0.7)

    card = add_rect(sl, x, 2.3, col_w, 8.55, fill_rgb=WHITE)

    # Mode image (webp will be converted)
    img_shape = add_pic(sl, img_name, x, 2.3, col_w, 3.2)
    if img_shape is None:
        # Fallback: color box with mode name
        color_box = add_rect(sl, x, 2.3, col_w, 3.2, fill_rgb=color)
        add_text(sl, x+0.2, 3.1, col_w-0.4, 1.4,
                 title.split()[-1], size=34, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Title bar
    title_bar = add_rect(sl, x, 5.5, col_w, 0.85, fill_rgb=color)
    add_text(sl, x+0.15, 5.6, col_w-0.3, 0.7, title,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(sl, x+0.25, 6.5, col_w-0.5, 0.55, subtitle,
             size=12, italic=True, color=MGRAY)

    # Bullets
    add_lines(sl, x+0.25, 7.15, col_w-0.5, 3.5,
        [(b, 13, False, DARK, PP_ALIGN.LEFT) for b in bullets])

send_back(sl, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 10 — Design Compacto & Sem Fio                   ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 10: Design & Wireless...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)
img = add_pic(sl, 'simse_hl_wireless.jpg', 10.5, 0, 9.5, 11.25)
dim = add_rect(sl, 10.5, 0, 9.5, 11.25, fill_rgb=NAVY2)

add_text(sl, 1, 0.9, 9, 0.7, 'DESIGN & CONECTIVIDADE',
         size=18, bold=True, color=BLUE)
add_text(sl, 1, 1.6, 9, 2.0, 'SEM FIO.\nSEM LIMITES.',
         size=54, bold=True, color=WHITE)
add_text(sl, 1, 4.0, 9, 1.0,
         'Liberdade total de movimento em qualquer ambiente industrial.',
         size=16, color=EEE)

feats = [
    (GREEN, 'Wi-Fi 6 Nativo',               '802.11a/b/g/n/ac — sem roteador externo necessário'),
    (BLUE,  'Baterias Hot-Swap',             'Troca a quente sem interromper a digitalização'),
    (GREEN, 'Modo Híbrido',                  'Cabeado (USB 3.0) ou sem fio — alterne quando precisar'),
    (BLUE,  'Computação de Borda',           'Processamento embarcado — não depende do PC em campo'),
    (GREEN, '203 × 80 × 44 mm · 600 g',    'Design de uma mão — scanner full-metal compacto'),
    (BLUE,  'Ambientes Hostis',              'Temperatura –10 °C a 40 °C · Laser Classe II (olho-seguro)'),
]
y = 5.4
for accent, title, desc in feats:
    add_rect(sl, 1.0, y, 0.25, 0.55, fill_rgb=accent)
    add_text(sl, 1.4, y,       7.5, 0.55, title, size=15, bold=True, color=WHITE)
    add_text(sl, 1.4, y+0.57,  7.5, 0.42, desc,  size=12, color=MGRAY)
    y += 1.08

send_back(sl, dim, img, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 11 — Geomagic Design X: Visão Geral              ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 11: Geomagic Design X — Visão Geral...")
sl = NS()

bg   = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)
acc  = add_rect(sl, 0, 0, 0.55, 11.25, fill_rgb=BLUE)
img  = add_pic(sl, 'simse_software_bg.webp', 10, 0, 10, 11.25)
dim  = add_rect(sl, 10, 0, 10, 11.25, fill_rgb=NAVY2)

add_text(sl, 1.0, 0.8, 8.5, 0.7, 'PARCEIRO DE SOFTWARE',
         size=15, bold=True, color=BLUE)
add_text(sl, 1.0, 1.5, 8.5, 1.5, 'Geomagic\nDesign X',
         size=50, bold=True, color=WHITE)
add_text(sl, 1.0, 3.5, 8.5, 0.7,
         'Software de Engenharia Reversa #1 do Mundo',
         size=18, bold=True, color=GREEN)
add_text(sl, 1.0, 4.4, 8.5, 2.2,
         'O Geomagic Design X transforma dados de escaneamento 3D em modelos '
         'CAD paramétricos completos — prontos para edição em SolidWorks, CATIA, '
         'Siemens NX, PTC Creo, Solid Edge e Inventor.',
         size=15, color=EEE)

caps = [
    (BLUE,  'Exact Surface',    'NURBS a partir de qualquer malha de scan'),
    (GREEN, 'LiveTransfer',     'Export direto e paramétrico para seu CAD'),
    (BLUE,  'Edição de Malha',  'Reparação automática de malhas de scan'),
    (GREEN, 'Análise de Desvio','Comparação CAD vs. escaneado em 3D'),
]
y = 7.0
for accent, title, desc in caps:
    dot = add_rect(sl, 1.0, y+0.04, 0.22, 0.42, fill_rgb=accent)
    add_text(sl, 1.35, y,      7, 0.5, title, size=15, bold=True, color=WHITE)
    add_text(sl, 1.35, y+0.5, 7, 0.45, desc, size=12, color=MGRAY)
    y += 1.0

send_back(sl, dim, img, acc, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 12 — Geomagic Design X: Funcionalidades          ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 12: Geomagic Design X — Funcionalidades...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=LGRAY)

add_text(sl, 0.5, 0.4, 19, 0.9, 'GEOMAGIC DESIGN X — FUNCIONALIDADES PRINCIPAIS',
         size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, 0.5, 1.35, 19, 0.55,
         'Fluxo completo: do ponto de escaneamento ao modelo paramétrico CAD',
         size=14, color=MGRAY, align=PP_ALIGN.CENTER)

feats = [
    (BLUE,  'Exact Surface',
     'Cria superfícies NURBS de alta qualidade a partir de malhas de scan, '
     'com controle total das curvaturas e continuidade de tangência.'),
    (GREEN, 'LiveTransfer®',
     'Transferência direta de dados paramétricos para SolidWorks, CATIA V5, '
     'Siemens NX, PTC Creo, Solid Edge e Inventor — editável como feature.'),
    (BLUE,  'Feature Recognition',
     'Reconhece automaticamente planos, cilindros, cones, esferas e toros '
     'da malha — reduz tempo de modelagem em até 80%.'),
    (GREEN, 'Análise de Desvio',
     'Compara o objeto escaneado com o CAD de referência com mapa de cores '
     '— ideal para inspeção e controle de qualidade dimensional.'),
    (BLUE,  'Edição de Malha',
     'Preenchimento automático de furos, suavização, remoção de ruído '
     'e alinhamento de múltiplas capturas em um único fluxo.'),
    (GREEN, 'Sketching no Scan',
     'Crie esboços 2D diretamente sobre a superfície do scan para '
     'extrusões, revoluções e cortes paramétricos precisos.'),
]

col_w = 5.8
col_h = 3.7
for i, (color, title, desc) in enumerate(feats):
    col = i % 3
    row = i // 3
    x = 0.5 + col * (col_w + 0.65)
    y = 2.2 + row * (col_h + 0.15)

    card = add_rect(sl, x, y, col_w, col_h, fill_rgb=WHITE)
    bar  = add_rect(sl, x, y, col_w, 0.55, fill_rgb=color)
    add_text(sl, x+0.2, y+0.08, col_w-0.4, 0.45, title,
             size=14, bold=True, color=WHITE)
    add_text(sl, x+0.2, y+0.7, col_w-0.4, 2.8, desc,
             size=13, color=DARK, wrap=True)

send_back(sl, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 13 — Fluxo Integrado: SIMSCAN-E → Geomagic       ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 13: Fluxo SIMSCAN-E → Geomagic...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)

add_text(sl, 0.5, 0.6, 19, 0.9, 'DO ESCANEAMENTO AO CAD PARAMÉTRICO',
         size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(sl, 0.5, 1.6, 19, 0.6,
         'Fluxo integrado SIMSCAN-E + Geomagic Design X — do objeto físico ao modelo digital editável',
         size=15, color=EEE, align=PP_ALIGN.CENTER)

steps = [
    (BLUE,  '01', 'ESCANEAR',   'SIMSCAN-E',
     'Captura a peça física em 3D com 0,020 mm de precisão e '
     '6,3M medições/segundo via Wi-Fi 6'),
    (GREEN, '02', 'PROCESSAR',  'DefinSight',
     'Malha 3D processada no software nativo — alinhamento, '
     'refinamento e exportação em .stl/.obj/.ply'),
    (BLUE,  '03', 'IMPORTAR',   'Geomagic Design X',
     'Importa a malha e converte em superfícies NURBS '
     'com Exact Surface e Feature Recognition automático'),
    (GREEN, '04', 'EXPORTAR',   'CAD / ERP / QA',
     'Modelo paramétrico Live Transfer para SolidWorks, '
     'CATIA, NX, Creo — pronto para produção'),
]

box_w = 4.3
for i, (color, num, step, tool, desc) in enumerate(steps):
    x = 0.5 + i * (box_w + 0.25)
    box = add_rect(sl, x, 2.8, box_w, 7.6, fill_rgb=DBOX)
    top_bar = add_rect(sl, x, 2.8, box_w, 0.12, fill_rgb=color)

    add_text(sl, x+0.2, 3.0,  box_w-0.4, 1.2, num,
             size=52, bold=True, color=color, align=PP_ALIGN.LEFT)
    add_text(sl, x+0.2, 4.4,  box_w-0.4, 0.65, step,
             size=20, bold=True, color=WHITE)
    add_text(sl, x+0.2, 5.15, box_w-0.4, 0.65, tool,
             size=13, bold=True, color=color)
    add_text(sl, x+0.2, 5.9,  box_w-0.4, 4.0, desc,
             size=13, color=EEE, wrap=True)

    # Arrow (except last)
    if i < 3:
        ax = x + box_w + 0.01
        add_rect(sl, ax, 5.8, 0.22, 0.5, fill_rgb=color)

send_back(sl, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 14 — Aplicações Industriais                      ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 14: Aplicações...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=WHITE)
accent_bar = add_rect(sl, 0, 0, 20, 1.45, fill_rgb=NAVY)

add_text(sl, 0.5, 0.25, 19, 0.9, 'APLICAÇÕES INDUSTRIAIS',
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

apps = [
    ('simse_app_aeroespacial.png', 'AEROESPACIAL',
     'Fuselagens, estruturas, ferramental de montagem e peças críticas em titânio e CFRP'),
    ('simse_app_automotivo.png',   'AUTOMOTIVO',
     'Carrocerias, parabrisas, painéis, estampados e protótipos em plástico e metal'),
    ('simse_app_moldes.png',       'MOLDES E FERRAMENTARIA',
     'Engenharia reversa de moldes, postais de referência e ferramental de injeção'),
    ('simse_app_energia.png',      'ENERGIA E IND. PESADA',
     'Turbinas, rotores, pás, válvulas e componentes de grande porte em campo'),
    ('simse_app_ferroviario.png',  'FERROVIÁRIO E NAVAL',
     'Chassi de vagões, estruturas de cascos, peças de grandes dimensões'),
    ('simse_app_outros.png',       'OUTROS SETORES',
     'Biomédico, odontológico, bens de consumo, preservação e design de produto'),
]

col_w = 5.9
col_h = 4.8
for i, (img_name, title, desc) in enumerate(apps):
    col = i % 3
    row = i // 3
    x   = 0.35 + col * (col_w + 0.52)
    y   = 1.6  + row * (col_h + 0.2)

    card = add_rect(sl, x, y, col_w, col_h, fill_rgb=LGRAY)
    img  = add_pic(sl, img_name, x, y, col_w, 3.0)

    bar = add_rect(sl, x, y + 3.0, col_w, 0.65, fill_rgb=NAVY)
    add_text(sl, x+0.15, y+3.07, col_w-0.3, 0.55, title,
             size=13, bold=True, color=WHITE)
    add_text(sl, x+0.15, y+3.8,  col_w-0.3, 0.9, desc,
             size=11, color=DARK, wrap=True)

send_back(sl, accent_bar, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 15 — Especificações Técnicas                     ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 15: Especificações Técnicas...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=LGRAY)
add_text(sl, 0.5, 0.35, 19, 0.9, 'ESPECIFICAÇÕES TÉCNICAS — SIMSCAN-E',
         size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Table simulation using colored rectangles + text
headers = ['PARÂMETRO', 'ESPECIFICAÇÃO', 'PARÂMETRO', 'ESPECIFICAÇÃO']
col_xs  = [0.4, 4.2, 10.4, 14.2]
col_ws  = [3.7, 5.9, 3.7, 5.4]

# Header row
for j, (h, cx, cw) in enumerate(zip(headers, col_xs, col_ws)):
    hr = add_rect(sl, cx, 1.45, cw, 0.6, fill_rgb=NAVY)
    add_text(sl, cx+0.15, 1.5, cw-0.3, 0.5, h,
             size=12, bold=True, color=WHITE)

specs_left = [
    ('Precisão',              '0,020 mm ¹'),
    ('Resolução',             '0,020 mm'),
    ('Velocidade de Medição', '6.300.000 pts/s'),
    ('Taxa de Quadros',       '180 FPS'),
    ('Linhas de Laser',       '81 (Triple Cross)'),
    ('Modo Ultrarrápido',     '63 linhas'),
    ('Modo Hyperfine',        '17 linhas'),
    ('Modo Furos Profundos',  '1 linha adicional'),
    ('Stand-off',             '300 mm'),
    ('Profundidade de Campo', '550 mm'),
]
specs_right = [
    ('Área de Varredura',     '700 × 600 mm'),
    ('Classe de Laser',       'Classe II (olho-seguro)'),
    ('Conectividade',         'Wi-Fi 6 + USB 3.0'),
    ('Dimensões',             '203 × 80 × 44 mm'),
    ('Peso',                  '600 g'),
    ('Bateria',               'Dupla hot-swap integrada'),
    ('Temperatura de Operação','–10 °C a 40 °C'),
    ('Formatos de Saída',     '.stl .obj .ply .asc .igs'),
    ('Software',              'DefinSight (nativo)'),
    ('Certificação',          'ISO 17025 ¹ ²'),
]

for row_i, ((p1, v1), (p2, v2)) in enumerate(zip(specs_left, specs_right)):
    ry     = 2.12 + row_i * 0.78
    bg_row = LGRAY if row_i % 2 == 0 else WHITE

    add_rect(sl, 0.4,  ry, 3.7, 0.72, fill_rgb=bg_row)
    add_rect(sl, 4.2,  ry, 5.9, 0.72, fill_rgb=bg_row)
    add_rect(sl, 10.4, ry, 3.7, 0.72, fill_rgb=bg_row)
    add_rect(sl, 14.2, ry, 5.4, 0.72, fill_rgb=bg_row)

    add_text(sl, 0.55,  ry+0.12, 3.5, 0.55, p1, size=12, bold=True,  color=DARK)
    add_text(sl, 4.35,  ry+0.12, 5.7, 0.55, v1, size=12, bold=False, color=NAVY)
    add_text(sl, 10.55, ry+0.12, 3.5, 0.55, p2, size=12, bold=True,  color=DARK)
    add_text(sl, 14.35, ry+0.12, 5.2, 0.55, v2, size=12, bold=False, color=NAVY)

# Footnotes
add_text(sl, 0.5, 10.4, 19, 0.8,
         '¹ Acreditado ISO 17025: PS(E) avaliado conforme VDI/VDE 2634 Pt.3 e JJF 1951.   '
         '² SD avaliado conforme VDI/VDE 2634 Pt.3 e JJF 1951.',
         size=10, color=MGRAY)

send_back(sl, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 16 — Suporte Técnico                             ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 16: Suporte Técnico...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)
img = add_pic(sl, 'simse_feat1_bg.jpg', 10, 0, 10, 11.25)
dim = add_rect(sl, 10, 0, 10, 11.25, fill_rgb=DBOX)

add_text(sl, 1, 0.9, 8.5, 0.7, 'PÓS-VENDA & SUPORTE',
         size=17, bold=True, color=BLUE)
add_text(sl, 1, 1.6, 8.5, 2.0, 'SUPORTE\nTÉCNICO',
         size=60, bold=True, color=WHITE)
add_text(sl, 1, 4.0, 8.5, 1.0,
         'Nossa equipe está disponível para garantir que sua operação nunca pare.',
         size=15, color=EEE)

supports = [
    (GREEN, 'Treinamento Completo',     'Implantação, operação e análise de resultados com o SIMSCAN-E'),
    (BLUE,  'Suporte Remoto',           'Atendimento técnico via videoconferência em até 4 horas úteis'),
    (GREEN, 'Calibração e Certificação','Serviço de calibração rastreável ao INMETRO quando necessário'),
    (BLUE,  'Assistência On-Site',      'Visita técnica para casos complexos ou treinamentos avançados'),
    (GREEN, 'Atualizações de Software', 'Licença DefinSight com atualizações incluídas no período de garantia'),
]
y = 5.3
for accent, title, desc in supports:
    add_rect(sl, 1.0, y, 0.22, 0.5, fill_rgb=accent)
    add_text(sl, 1.38, y,       7.4, 0.5, title, size=15, bold=True, color=WHITE)
    add_text(sl, 1.38, y+0.52,  7.4, 0.45, desc, size=12, color=MGRAY)
    y += 1.04

send_back(sl, dim, img, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 17 — Proposta Comercial                          ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 17: Proposta Comercial...")
sl = NS()

bg = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)

add_text(sl, 0.5, 0.7, 19, 0.9, 'PROPOSTA COMERCIAL',
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 0.5, 1.65, 19, 0.7,
         'Sistema SIMSCAN-E + Geomagic Design X — Enterfix',
         size=18, color=GREEN, align=PP_ALIGN.CENTER)

# Package cards
packages = [
    (NAVY2, BLUE,  'SIMSCAN-E',
     'Scanner 3D Industrial',
     ['✓ Hardware SIMSCAN-E completo',
      '✓ Software DefinSight (licença nativa)',
      '✓ Maleta de transporte premium',
      '✓ Acessórios e kit de calibração',
      '✓ Treinamento de operação (2 dias)',
      '✓ Suporte técnico 12 meses'],
     'Consultar'),
    (DBOX, GREEN, 'GEOMAGIC DESIGN X',
     'Engenharia Reversa CAD',
     ['✓ Licença Geomagic Design X',
      '✓ LiveTransfer para SolidWorks / CATIA / NX',
      '✓ Treinamento de uso (2 dias)',
      '✓ 1 ano de manutenção STI incluída',
      '✓ Suporte via helpdesk 3D Systems',
      '✓ Material de treinamento digital'],
     'Consultar'),
    (RGBColor(0x1A, 0x3A, 0x1A), GREEN, 'COMBO COMPLETO',
     'SIMSCAN-E + Geomagic Design X',
     ['✓ Tudo do SIMSCAN-E +',
      '✓ Tudo do Geomagic Design X +',
      '✓ Treinamento integrado scan→CAD',
      '✓ Projeto piloto com peça do cliente',
      '✓ Suporte prioritário 12 meses',
      '✓ Melhor custo-benefício'],
     'Consultar'),
]

card_w = 5.8
for i, (bg_c, accent, title, sub, items, price) in enumerate(packages):
    x = 0.5 + i * (card_w + 0.65)
    card = add_rect(sl, x, 2.6, card_w, 8.1, fill_rgb=bg_c)
    bar  = add_rect(sl, x, 2.6, card_w, 0.12, fill_rgb=accent)

    add_text(sl, x+0.2, 2.8,  card_w-0.4, 0.8, title,
             size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, x+0.2, 3.65, card_w-0.4, 0.65, sub,
             size=13, color=accent, align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_text(sl, x+0.25, 4.5 + j*0.88, card_w-0.5, 0.75,
                 item, size=13, color=EEE)

    add_rect(sl, x+0.4, 9.7, card_w-0.8, 0.7, fill_rgb=accent)
    add_text(sl, x+0.4, 9.78, card_w-0.8, 0.58, f'💬 {price}',
             size=14, bold=True, color=WHITE if accent != GREEN else NAVY,
             align=PP_ALIGN.CENTER)

send_back(sl, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 18 — Vantagens de Ser Nosso Cliente              ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 18: Vantagens...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=LGRAY)
hdr = add_rect(sl, 0, 0, 20, 1.55, fill_rgb=NAVY)

add_text(sl, 0.5, 0.28, 19, 1.0, 'VANTAGENS DE SER NOSSO CLIENTE',
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

advantages = [
    (BLUE,  '🏆', 'Equipe Técnica Especializada',
     'Profissionais certificados em metrologia 3D e engenharia reversa'),
    (GREEN, '⚡', 'Rapidez no Atendimento',
     'Resposta em até 4 horas úteis — suporte remoto ou presencial'),
    (BLUE,  '🎯', 'Produtos de Alta Qualidade',
     'Hardware e software das marcas mais confiáveis do mercado global'),
    (GREEN, '🤝', 'Parceria de Longo Prazo',
     'Acompanhamos desde a implantação até a evolução dos processos'),
    (BLUE,  '💡', 'Treinamento e Capacitação',
     'Treinamentos técnicos personalizados para sua equipe e aplicação'),
    (GREEN, '🌐', 'Suporte Web e On-site',
     'Atendimento remoto, presencial e via WebEx / Teams / WhatsApp'),
]

card_w = 5.8
card_h = 3.9
for i, (color, icon, title, desc) in enumerate(advantages):
    col = i % 3
    row = i // 3
    x   = 0.5 + col * (card_w + 0.65)
    y   = 1.8 + row * (card_h + 0.25)

    card = add_rect(sl, x, y, card_w, card_h, fill_rgb=WHITE)
    add_text(sl, x+0.2, y+0.3,  card_w-0.4, 0.9, icon,
             size=36, align=PP_ALIGN.CENTER)
    bar = add_rect(sl, x, y+1.25, card_w, 0.08, fill_rgb=color)
    add_text(sl, x+0.2, y+1.45, card_w-0.4, 0.75, title,
             size=15, bold=True, color=NAVY)
    add_text(sl, x+0.2, y+2.25, card_w-0.4, 1.5, desc,
             size=13, color=MGRAY, wrap=True)

send_back(sl, hdr, bg)

# ╔══════════════════════════════════════════════════════════╗
# ║  SLIDE 19 — CTA Final / Contato                         ║
# ╚══════════════════════════════════════════════════════════╝
print("Construindo slide 19: CTA Final...")
sl = NS()

bg  = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY)
img = add_pic(sl, 'simse_hl_dimensoes.jpg', 0, 0, 20, 11.25)
dim = add_rect(sl, 0, 0, 20, 11.25, fill_rgb=NAVY2)

# Centered content box
box = add_rect(sl, 4.0, 2.0, 12, 7.25, fill_rgb=DBOX)
bar = add_rect(sl, 4.0, 2.0, 12, 0.14, fill_rgb=GREEN)

add_text(sl, 4.2, 2.3, 11.6, 0.85,
         'PRONTO PARA TRANSFORMAR SEU PROCESSO?',
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, 4.2, 3.25, 11.6, 1.5,
         'Solicite uma demonstração do SIMSCAN-E com peças reais da sua linha de produção '
         'e descubra o que 6,3 milhões de medições por segundo fazem pela sua qualidade.',
         size=15, color=EEE, align=PP_ALIGN.CENTER)

add_text(sl, 4.2, 5.0, 11.6, 0.65,
         'Enterfix Indústria, Comércio e Serviços Ltda.',
         size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

contacts = [
    ('📞  (11) 4942-2222   |   WhatsApp: (11) 96841-6776', WHITE, 14),
    ('✉️  vendas@enterfix.com.br   |   www.enterfix.com.br', EEE,  13),
    ('📍  Rua Waldemar Martins Ferreira, 287 — São Bernardo do Campo - SP', MGRAY, 12),
    ('🕐  Seg–Sex 08:00–12:00 e 13:00–17:30', MGRAY, 12),
]
cy = 5.9
for text, color, size in contacts:
    add_text(sl, 4.2, cy, 11.6, 0.6, text,
             size=size, color=color, align=PP_ALIGN.CENTER)
    cy += 0.65

send_back(sl, dim, img, bg)

# ─────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────
print(f"\nSalvando: {DST}")
prs.save(DST)

total_slides = len(list(prs.slides))
print(f"\n✅ Apresentação salva com sucesso!")
print(f"   Total de slides: {total_slides}")
print(f"   Slides 1-6:  Empresa Enterfix (originais)")
print(f"   Slides 7-19: SIMSCAN-E + Geomagic Design X (novos)")
print(f"\n   Arquivo: {DST}")
